from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Paleta oscura (consistente con figuras del paper)
BG        = RGBColor(0x0F, 0x17, 0x2A)   # fondo principal
TEXT_PRI  = RGBColor(0xF1, 0xF5, 0xF9)   # texto principal
TEXT_SEC  = RGBColor(0xCB, 0xD5, 0xE1)   # texto secundario
TEXT_MUTE = RGBColor(0x94, 0xA3, 0xB8)   # texto muted
BLUE      = RGBColor(0x60, 0xA5, 0xFA)   # azul acento
GREEN     = RGBColor(0x4A, 0xDE, 0x80)   # verde acento
YELLOW    = RGBColor(0xFB, 0xBF, 0x24)   # amarillo acento
ORANGE    = RGBColor(0xFB, 0x92, 0x3C)   # naranja acento
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)   # morado acento
RED_BG    = RGBColor(0x7F, 0x1D, 0x1D)   # fondo rojo (negativo)
GREEN_BG  = RGBColor(0x14, 0x53, 0x2D)   # fondo verde (positivo)
DARK_BOX  = RGBColor(0x1E, 0x29, 0x3B)   # fondo de cuadros

# Dimensiones (16:9 estándar)
W = Inches(13.33)
H = Inches(7.5)


def set_background(slide, color=BG):
    """Aplica color sólido de fondo a un slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=24, color=TEXT_PRI, bold=False,
                align=PP_ALIGN.LEFT, italic=False):
    """Agrega un cuadro de texto y retorna el shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_title(slide, text, color=TEXT_PRI, size=36):
    """Título del slide en la banda superior."""
    return add_textbox(
        slide, text,
        left=Inches(0.5), top=Inches(0.3),
        width=Inches(12.3), height=Inches(0.9),
        font_size=size, color=color, bold=True,
        align=PP_ALIGN.LEFT
    )


def add_image(slide, path, left, top, width, height=None):
    """Inserta imagen; si height=None la escala por ancho."""
    if height is not None:
        return slide.shapes.add_picture(str(path), left, top, width, height)
    return slide.shapes.add_picture(str(path), left, top, width)


def add_speaker_note(slide, text):
    """Agrega guión en las notas del presentador."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_pt=0):
    """Agrega un rectángulo de color sólido."""
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape
