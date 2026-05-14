"""
Genera MTE_Avances_2026.pptx — Revisión de avances MTE, mayo 2026.
Ejecutar: python presentation/build_presentation.py
Output:   presentation/MTE_Avances_2026.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from utils import (
    set_background, add_textbox, add_title, add_image,
    add_speaker_note, add_rect,
    BG, TEXT_PRI, TEXT_SEC, TEXT_MUTE,
    BLUE, GREEN, YELLOW, ORANGE, PURPLE, RED_BG, GREEN_BG, DARK_BOX,
    W, H
)

FIGURES = Path(__file__).parent.parent / "outputs" / "paper"


def new_blank_slide(prs):
    """Agrega slide en blanco con fondo oscuro."""
    layout = prs.slide_layouts[6]   # blank layout
    slide = prs.slides.add_slide(layout)
    set_background(slide)
    return slide


# ──────────────────────────────────────────────
# BLOQUE 1 — GANCHO (slides 1–3)
# ──────────────────────────────────────────────

def slide_01_portada(prs):
    """Portada."""
    s = new_blank_slide(prs)
    # Línea decorativa superior
    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), BLUE)
    # Título principal
    add_textbox(
        s, "Mercados P2P de Energía:\n¿Qué nos dicen los datos\nreales de Pasto?",
        Inches(1), Inches(1.2), Inches(11.3), Inches(2.8),
        font_size=40, color=TEXT_PRI, bold=True, align=PP_ALIGN.CENTER
    )
    # Autor
    add_textbox(
        s, "Brayan S. López-Mendez",
        Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
        font_size=22, color=BLUE, bold=True, align=PP_ALIGN.CENTER
    )
    # Institución y asesores
    add_textbox(
        s, "Maestría en Ingeniería Electrónica · Universidad de Nariño\n"
           "Asesores: Andrés Pantoja · Germán Obando",
        Inches(1), Inches(4.8), Inches(11.3), Inches(0.8),
        font_size=16, color=TEXT_MUTE, align=PP_ALIGN.CENTER
    )
    # Evento y fecha
    add_textbox(
        s, "Revisión de avances proyecto MTE · Mayo 2026",
        Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
        font_size=14, color=TEXT_MUTE, italic=True, align=PP_ALIGN.CENTER
    )
    # Línea decorativa inferior
    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), BLUE)
    add_speaker_note(s,
        "Buenos días. Hoy les voy a contar qué pasa cuando cinco instituciones "
        "de Pasto comparten su energía solar — y qué nos dice eso sobre las "
        "reglas que rigen la energía en Colombia."
    )
    return s


def slide_02_pregunta_provocadora(prs):
    """La pregunta provocadora."""
    s = new_blank_slide(prs)
    add_title(s, "¿La regulación colombiana está lista para el solar de 2030?", color=TEXT_PRI)
    # Cita UPME
    add_rect(s, Inches(1.5), Inches(1.4), Inches(10.3), Inches(1.6), DARK_BOX)
    add_textbox(
        s, '"Colombia instalará 5× más solar entre 2025 y 2030."',
        Inches(1.7), Inches(1.5), Inches(9.9), Inches(0.9),
        font_size=24, color=YELLOW, bold=True, align=PP_ALIGN.CENTER
    )
    add_textbox(
        s, "— Plan de Expansión UPME 2025–2039",
        Inches(1.7), Inches(2.3), Inches(9.9), Inches(0.5),
        font_size=14, color=TEXT_MUTE, italic=True, align=PP_ALIGN.CENTER
    )
    # Pregunta central
    add_textbox(
        s, "¿Los mecanismos de liquidación de la CREG 174 y CREG 101 072\n"
           "están diseñados para aprovechar esa energía?",
        Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.2),
        font_size=22, color=TEXT_PRI, align=PP_ALIGN.CENTER
    )
    # Respuesta anticipada
    add_rect(s, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.7), GREEN_BG)
    add_textbox(
        s, "Esta presentación responde esa pregunta con datos reales de Pasto.",
        Inches(2.7), Inches(4.85), Inches(7.9), Inches(0.6),
        font_size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "El Plan de Expansión proyecta 5 veces más paneles solares en 2030. Eso es "
        "enorme. Pero nadie ha respondido con datos reales de acá: ¿la regulación "
        "que tenemos hoy está diseñada para aprovechar esa energía? ¿Cuánto dinero "
        "se deja sobre la mesa?"
    )
    return s


def slide_03_mte_intro(prs):
    """El proyecto MTE — 5 instituciones."""
    s = new_blank_slide(prs)
    add_title(s, "5 instituciones · 744 horas · datos reales de Pasto", color=BLUE)
    # 5 instituciones en grid
    insts = [
        ("🏫 Udenar", "Universidad de Nariño"),
        ("🏫 U. Mariana", "Universidad Mariana"),
        ("🏫 UCC", "U. Cooperativa de Colombia"),
        ("🏥 HUDN", "Hospital Universitario Dpto. Nariño"),
        ("🏫 Cesmag", "Universidad Cesmag"),
    ]
    cols = [Inches(0.5), Inches(3.2), Inches(5.9), Inches(8.6)]  # 4 posiciones fila superior
    for i, (icon_name, full) in enumerate(insts):
        x = cols[i] if i < 4 else Inches(5.9)   # Cesmag: centrado en fila 2
        y = Inches(1.5) if i < 4 else Inches(3.8)
        add_rect(s, x, y, Inches(2.5), Inches(1.6), DARK_BOX)
        add_textbox(s, icon_name, x + Inches(0.1), y + Inches(0.1),
                    Inches(2.3), Inches(0.6), font_size=18, color=BLUE, bold=True)
        add_textbox(s, full, x + Inches(0.1), y + Inches(0.65),
                    Inches(2.3), Inches(0.8), font_size=11, color=TEXT_SEC)
    # Stats
    add_textbox(
        s, "⚡  9.9 kW instalados     📊  1 medición cada 2 min     📅  Agosto 2025 completo",
        Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.6),
        font_size=16, color=GREEN, align=PP_ALIGN.CENTER
    )
    add_textbox(
        s, "Datos medidos en campo — no estimaciones, no promedios nacionales.",
        Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5),
        font_size=13, color=TEXT_MUTE, italic=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "El proyecto MTE lo conforman 5 instituciones de Pasto: Udenar, Mariana, UCC, "
        "el Hospital Universitario y Cesmag. En agosto de 2025 instalamos sensores en "
        "cada una y durante un mes completo registramos cada dos minutos cuánta energía "
        "se consumía y cuánta generaban los paneles. Eso nos da 744 horas de datos "
        "reales — no encuestas, no estimaciones nacionales, datos medidos acá en Pasto."
    )
    return s


# ──────────────────────────────────────────────
# BLOQUE 2 — PREGUNTA 1 (slides 4–6)
# ──────────────────────────────────────────────

def slide_04_que_es_mte(prs):
    """¿Qué mide el MTE?"""
    s = new_blank_slide(prs)
    add_title(s, "El MTE mide la energía en tiempo real", color=BLUE)
    items = [
        ("⚡", "9.9 kW", "de paneles solares instalados en 5 instituciones"),
        ("📊", "Cada 2 min", "datos registrados → agrupados por hora para el análisis"),
        ("📅", "744 horas", "agosto 2025 completo — mes de mayor irradiancia y actividad académica"),
        ("🏙️", "Pasto, Nariño", "datos locales, no nacionales ni interpolados"),
    ]
    for i, (icon, bold_text, rest) in enumerate(items):
        y = Inches(1.5) + i * Inches(1.15)
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.95), DARK_BOX)
        add_textbox(s, icon, Inches(0.7), y + Inches(0.2), Inches(0.5), Inches(0.6), font_size=22)
        add_textbox(s, bold_text, Inches(1.4), y + Inches(0.15), Inches(2.2), Inches(0.6),
                    font_size=20, color=GREEN, bold=True)
        add_textbox(s, rest, Inches(3.7), y + Inches(0.2), Inches(8.8), Inches(0.6),
                    font_size=16, color=TEXT_SEC)
    add_speaker_note(s,
        "Estos cuatro indicadores resumen qué tan grande fue la medición. "
        "9.9 kilovatios de paneles, cinco instituciones, una lectura cada dos minutos "
        "durante agosto completo. Eso da 744 horas. ¿Por qué agosto? Es el mes de mayor "
        "irradiancia solar en Pasto y con las cinco instituciones en plena actividad. "
        "Los datos se agruparon por hora para el análisis — sin perder la resolución."
    )
    return s


def slide_05_perfiles(prs):
    """Dos perfiles, una oportunidad — figura."""
    s = new_blank_slide(prs)
    add_title(s, "Dos instituciones, dos realidades — la oportunidad de intercambio", color=TEXT_PRI)
    fig = FIGURES / "fig_paper_profiles_2agents.png"
    add_image(s, fig, Inches(0.4), Inches(1.1), Inches(8.5))
    add_rect(s, Inches(9.1), Inches(1.2), Inches(3.9), Inches(2.2), DARK_BOX)
    add_textbox(s, "🏥 HUDN (Hospital)",
                Inches(9.3), Inches(1.3), Inches(3.5), Inches(0.5),
                font_size=14, color=BLUE, bold=True)
    add_textbox(s, "Consumo constante 24/7\nSolar variable de día\n→ Siempre necesita energía",
                Inches(9.3), Inches(1.8), Inches(3.5), Inches(1.4),
                font_size=12, color=TEXT_SEC)
    add_rect(s, Inches(9.1), Inches(3.7), Inches(3.9), Inches(2.2), DARK_BOX)
    add_textbox(s, "🏫 Udenar",
                Inches(9.3), Inches(3.8), Inches(3.5), Inches(0.5),
                font_size=14, color=GREEN, bold=True)
    add_textbox(s, "Consumo variable (clases)\nSolar fuerte al mediodía\n→ A veces le sobra energía",
                Inches(9.3), Inches(4.3), Inches(3.5), Inches(1.4),
                font_size=12, color=TEXT_SEC)
    add_textbox(
        s, "Cuando a Udenar le sobra → el hospital la necesita.",
        Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.6),
        font_size=16, color=YELLOW, bold=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "Si miran estas dos instituciones — el hospital y la universidad — ven algo "
        "interesante. El hospital consume energía todo el día, a toda hora: urgencias, "
        "UCI, climatización. La universidad, en cambio, tiene picos cuando hay clases "
        "y al mediodía le sobra solar. Esa diferencia es exactamente la oportunidad: "
        "cuando a Udenar le sobra energía, el HUDN la necesita. La pregunta es: "
        "¿cómo organizamos eso?"
    )
    return s


def slide_06_problema_fondo(prs):
    """El problema de fondo: mensual vs horario."""
    s = new_blank_slide(prs)
    add_title(s, "¿Qué hace hoy la regulación con la energía que sobra?", color=TEXT_PRI)
    add_rect(s, Inches(0.5), Inches(1.3), Inches(5.9), Inches(4.2), RED_BG)
    add_textbox(s, "Regulación actual (C1 / C2)",
                Inches(0.7), Inches(1.45), Inches(5.5), Inches(0.6),
                font_size=18, color=RGBColor(0xFC, 0xA5, 0xA5), bold=True)
    add_textbox(
        s, "Acumula toda la energía sobrante\ndel mes en un 'tanque'.\n\n"
           "Al final del mes → liquida todo\nal precio promedio mensual.\n\n"
           "No distingue si sobró a mediodía\ncon sol pleno o a las 2am.",
        Inches(0.7), Inches(2.1), Inches(5.5), Inches(3.2),
        font_size=15, color=TEXT_PRI
    )
    add_rect(s, Inches(6.9), Inches(1.3), Inches(5.9), Inches(4.2), GREEN_BG)
    add_textbox(s, "Mercado P2P",
                Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.6),
                font_size=18, color=GREEN, bold=True)
    add_textbox(
        s, "En cada hora, quien sobra le\nvende directamente a quien\nnecesita.\n\n"
           "Precio acordado en ese momento\nsegún oferta y demanda real.\n\n"
           "Captura el valor de cada hora\nde forma independiente.",
        Inches(7.1), Inches(2.1), Inches(5.5), Inches(3.2),
        font_size=15, color=TEXT_PRI
    )
    add_textbox(s, "→", Inches(6.2), Inches(3.1), Inches(0.6), Inches(0.8),
                font_size=32, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        s, "¿El mecanismo horario da mejores resultados que el mensual?\nEso es la Pregunta 2.",
        Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.7),
        font_size=15, color=BLUE, italic=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "Hoy la regulación funciona así: toda la energía que no usas en el mes va a un "
        "'tanque común' y al final del mes te pagan por ella a un precio promedio. No "
        "importa si sobró a mediodía con sol pleno o a las dos de la mañana. Todo vale "
        "lo mismo. El P2P dice: 'no, en cada hora hagamos una transacción directa'. "
        "¿Eso da mejores resultados? Esa es la Pregunta 2."
    )
    return s


# ──────────────────────────────────────────────
# BLOQUE 3 — PREGUNTA 2 (slides 7–13)
# ──────────────────────────────────────────────

def slide_07_que_es_p2p(prs):
    """¿Qué es un mercado P2P? Analogía."""
    s = new_blank_slide(prs)
    add_title(s, "P2P: comprar y vender energía directamente — sin intermediario", color=TEXT_PRI)
    add_rect(s, Inches(0.4), Inches(1.3), Inches(5.8), Inches(4.5), RED_BG)
    add_textbox(s, "Sin P2P — hoy",
                Inches(0.6), Inches(1.45), Inches(5.4), Inches(0.55),
                font_size=18, color=RGBColor(0xFC, 0xA5, 0xA5), bold=True)
    add_textbox(
        s, "Udenar\n    ↓\nEmpresa eléctrica\n    ↓\nHUDN\n\n"
           "El intermediario fija el precio\ny se queda con el margen.",
        Inches(0.6), Inches(2.1), Inches(5.4), Inches(3.4),
        font_size=16, color=TEXT_PRI
    )
    add_rect(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(4.5), GREEN_BG)
    add_textbox(s, "Con P2P",
                Inches(7.2), Inches(1.45), Inches(5.4), Inches(0.55),
                font_size=18, color=GREEN, bold=True)
    add_textbox(
        s, "Udenar  ↔  HUDN\n\nPrecio acordado\ndirectamente entre ellos.\n\n"
           "Como Airbnb o Rappi:\ntransacción persona a persona.",
        Inches(7.2), Inches(2.1), Inches(5.4), Inches(3.4),
        font_size=16, color=TEXT_PRI
    )
    add_textbox(s, "↔", Inches(6.2), Inches(3.1), Inches(0.6), Inches(0.8),
                font_size=32, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        s, "El algoritmo hace la negociación automáticamente, hora por hora.",
        Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.6),
        font_size=15, color=BLUE, italic=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "¿Qué es un mercado P2P? La idea es simple: en lugar de que la energía sobrante "
        "pase por la empresa eléctrica antes de llegar a quien la necesita, las dos "
        "instituciones negocian directamente. Como cuando alguien alquila su casa en "
        "Airbnb en lugar de ir a un hotel, o pide comida a un restaurante en Rappi "
        "directamente. El algoritmo hace toda la negociación automáticamente, "
        "hora por hora."
    )
    return s


def slide_08_algoritmo(prs):
    """Cómo funciona el algoritmo — sin ecuaciones."""
    s = new_blank_slide(prs)
    add_title(s, "El algoritmo: dos pasos que se repiten hasta encontrar precio justo", color=PURPLE)
    pasos = [
        ("1", PURPLE, "Vendedores anuncian precio",
         "Cada institución con energía sobrante dice cuánto cobra (COP/kWh), "
         "buscando su mayor ganancia dentro de lo que el mercado acepta."),
        ("2", BLUE, "Compradores responden",
         "Los que necesitan energía ajustan cuánto compran a cada vendedor "
         "según el precio — prefieren al más barato."),
        ("✓", GREEN, "Se alcanza el equilibrio",
         "Precios y cantidades convergen en milisegundos. Resultado: el precio "
         "de mercado justo para ese momento exacto."),
    ]
    for i, (num, col, titulo, desc) in enumerate(pasos):
        y = Inches(1.5) + i * Inches(1.5)
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(1.25), DARK_BOX)
        add_textbox(s, num, Inches(0.6), y + Inches(0.3), Inches(0.6), Inches(0.7),
                    font_size=26, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, titulo, Inches(1.4), y + Inches(0.1), Inches(4.5), Inches(0.55),
                    font_size=17, color=col, bold=True)
        add_textbox(s, desc, Inches(1.4), y + Inches(0.65), Inches(11.2), Inches(0.55),
                    font_size=13, color=TEXT_SEC)
    fig = FIGURES / "fig_paper_convergence_h0512.png"
    add_image(s, fig, Inches(0.4), Inches(6.0), Inches(12.5), Inches(1.25))
    add_speaker_note(s,
        "El algoritmo hace dos pasos. Primero, cada institución con energía sobrante "
        "anuncia su precio, buscando la mayor ganancia posible. Segundo, los que "
        "necesitan energía ajustan cuánto compran a cada vendedor — prefieren al más "
        "barato. Esos dos pasos se repiten hasta que todos los precios y cantidades "
        "se estabilizan. Eso pasa en milisegundos. La figura muestra cómo converge "
        "en una hora real de nuestros datos."
    )
    return s


def slide_09_esquemas_colombianos(prs):
    """Los dos esquemas regulatorios colombianos."""
    s = new_blank_slide(prs)
    add_title(s, "¿Con qué comparamos el P2P? Los dos esquemas colombianos", color=TEXT_PRI)
    add_rect(s, Inches(0.4), Inches(1.3), Inches(5.8), Inches(4.5), RGBColor(0x1E, 0x3A, 0x5F))
    add_textbox(s, "C1 · CREG 174/2021",
                Inches(0.6), Inches(1.45), Inches(5.4), Inches(0.55),
                font_size=18, color=BLUE, bold=True)
    add_textbox(s, "Autogeneración individual (AGPE)",
                Inches(0.6), Inches(2.05), Inches(5.4), Inches(0.45),
                font_size=13, color=TEXT_MUTE, italic=True)
    add_textbox(
        s, "• Cada institución gestiona\n  su propio excedente.\n\n"
           "• Liquidación mensual al\n  precio promedio del mes.\n\n"
           "• Es lo que existe HOY para\n  autogeneradores en Colombia.",
        Inches(0.6), Inches(2.55), Inches(5.4), Inches(3.0),
        font_size=15, color=TEXT_PRI
    )
    add_rect(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(4.5), RGBColor(0x1E, 0x3A, 0x2F))
    add_textbox(s, "C2 · CREG 101 072/2025",
                Inches(7.2), Inches(1.45), Inches(5.4), Inches(0.55),
                font_size=18, color=GREEN, bold=True)
    add_textbox(s, "Autogeneración colectiva (AGRC)",
                Inches(7.2), Inches(2.05), Inches(5.4), Inches(0.45),
                font_size=13, color=TEXT_MUTE, italic=True)
    add_textbox(
        s, "• Las 5 instituciones forman\n  una comunidad energética.\n\n"
           "• Excedentes compartidos según\n  capacidad instalada (PDE).\n\n"
           "• Es la NUEVA regulación para\n  comunidades (vigente 2025).",
        Inches(7.2), Inches(2.55), Inches(5.4), Inches(3.0),
        font_size=15, color=TEXT_PRI
    )
    add_textbox(
        s, "Los tres mecanismos usan exactamente los mismos datos MTE — solo cambia la regla de liquidación.",
        Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.6),
        font_size=14, color=TEXT_MUTE, italic=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "Para saber si el P2P es mejor, necesitamos compararlo con algo. Usamos dos "
        "esquemas que ya existen en Colombia. El C1 es la CREG 174, que es lo que "
        "tienen hoy los autogeneradores individuales: liquidan su excedente mensualmente. "
        "El C2 es la nueva regulación del 2025 para comunidades energéticas: las cinco "
        "instituciones forman un grupo y comparten sus excedentes según el tamaño de "
        "cada panel. Los tres mecanismos usan los mismos datos MTE — solo cambia la "
        "regla de liquidación."
    )
    return s


def slide_10_resultado_principal(prs):
    """Resultado principal — tabla en COP."""
    s = new_blank_slide(prs)
    add_title(s, "P2P genera más beneficio neto en agosto 2025 — φ = 1.5 (144% cobertura)", color=YELLOW)
    rows = [
        ("🥇 P2P", "4.12 M COP", "1.81 M COP", "5.94 M COP", "+5.5% vs C1", GREEN_BG, GREEN),
        ("🥈 C2 (Colectivo)", "4.12 M COP", "1.62 M COP", "5.75 M COP", "+2.1% vs C1",
         RGBColor(0x1E, 0x3A, 0x2F), RGBColor(0x86, 0xEF, 0xAC)),
        ("🥉 C1 (Individual)", "4.12 M COP", "1.50 M COP", "5.63 M COP", "referencia",
         DARK_BOX, TEXT_MUTE),
    ]
    headers = ["Escenario", "Autoconsumo", "Excedente", "TOTAL", "vs C1"]
    col_x = [Inches(0.4), Inches(3.6), Inches(5.9), Inches(8.2), Inches(10.5)]
    col_w = [Inches(3.0), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.5)]
    for j, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_textbox(s, hdr, x, Inches(1.3), w, Inches(0.45),
                    font_size=13, color=TEXT_MUTE, bold=True)
    for i, (scen, auto, exc, total, delta, bg, fg) in enumerate(rows):
        y = Inches(1.85) + i * Inches(1.3)
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(1.1), bg)
        vals = [scen, auto, exc, total, delta]
        for j, (val, x, w) in enumerate(zip(vals, col_x, col_w)):
            add_textbox(s, val, x + Inches(0.05), y + Inches(0.2), w - Inches(0.1), Inches(0.7),
                        font_size=15 if j == 0 else 16,
                        color=fg if j in (0, 3, 4) else TEXT_PRI,
                        bold=(j in (0, 3)))
    fig = FIGURES / "fig_paper_ahorro_decomposition.png"
    add_image(s, fig, Inches(0.4), Inches(5.2), Inches(12.5), Inches(1.9))
    add_speaker_note(s,
        "Aquí están los resultados. En agosto 2025, con la cobertura solar del caso "
        "de estudio, P2P genera 5.94 millones de pesos de beneficio neto para las "
        "cinco instituciones. C2 genera 5.75 y C1 genera 5.63 millones. "
        "P2P gana 5.5% sobre C1 (la regulación individual) y 3.3% sobre C2 "
        "(la colectiva); C2 a su vez gana 2.1% sobre C1. "
        "Importante: el autoconsumo — los 4.12 millones — es idéntico en los tres. "
        "La diferencia completa viene del excedente. La figura lo visualiza."
    )
    return s


def slide_11_por_agente(prs):
    """Por institución: quién prefiere qué."""
    s = new_blank_slide(prs)
    add_title(s, "4 de 5 instituciones prefieren P2P individualmente", color=TEXT_PRI)
    fig = FIGURES / "fig_paper_per_agent_benefit.png"
    add_image(s, fig, Inches(0.4), Inches(1.1), Inches(7.8))
    data = [
        ("Udenar",  "982 k",  "961 k",  "889 k",  "P2P ✓", GREEN),
        ("Mariana", "1,073 k","1,061 k","1,038 k","P2P ✓", GREEN),
        ("UCC",     "1,700 k","1,658 k","1,748 k","C2 ✓",  BLUE),
        ("HUDN",    "987 k",  "968 k",  "914 k",  "P2P ✓", GREEN),
        ("Cesmag",  "1,196 k","982 k",  "1,171 k","P2P ✓", GREEN),
    ]
    headers = ["Institución", "P2P", "C1", "C2", "Mejor"]
    col_x = [Inches(8.4), Inches(9.6), Inches(10.5), Inches(11.4), Inches(12.3)]
    col_w = [Inches(1.1), Inches(0.85), Inches(0.85), Inches(0.85), Inches(1.0)]
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_textbox(s, h, x, Inches(1.2), w, Inches(0.4),
                    font_size=11, color=TEXT_MUTE, bold=True)
    for i, (inst, p2p, c1, c2, best, col) in enumerate(data):
        y = Inches(1.7) + i * Inches(0.95)
        bg = DARK_BOX if i % 2 == 0 else RGBColor(0x0F, 0x17, 0x2A)
        add_rect(s, Inches(8.3), y, Inches(4.7), Inches(0.85), bg)
        for j, (val, x, w) in enumerate(zip([inst, p2p, c1, c2, best], col_x, col_w)):
            add_textbox(s, val, x, y + Inches(0.15), w, Inches(0.55),
                        font_size=11,
                        color=col if j in (0, 4) else TEXT_PRI,
                        bold=(j == 4))
    add_rect(s, Inches(8.3), Inches(6.55), Inches(4.7), Inches(0.8), RGBColor(0x1E, 0x3A, 0x5F))
    add_textbox(
        s, "UCC prefiere C2: es la institución con mayor consumo. "
           "El reparto colectivo le asigna más créditos que el P2P.",
        Inches(8.4), Inches(6.6), Inches(4.5), Inches(0.7),
        font_size=10, color=BLUE
    )
    add_speaker_note(s,
        "Cuatro de las cinco instituciones prefieren P2P. La excepción es UCC. "
        "¿Por qué UCC prefiere C2? UCC es la institución con mayor consumo de "
        "las cinco. Con C2, el reparto colectivo le asigna una proporción grande "
        "de los excedentes comunitarios en función de su capacidad instalada. "
        "Eso le da más créditos que lo que gana negociando hora a hora en P2P. "
        "Es un resultado interesante: la regulación colectiva puede favorecer a "
        "los grandes consumidores."
    )
    return s


def slide_12_transicion_fase(prs):
    """La transición de fase — slide estrella."""
    s = new_blank_slide(prs)
    add_title(s, "¿Desde cuánta generación solar conviene el P2P?", color=YELLOW)
    fig = FIGURES / "fig_paper_c1_vs_c4_detailed.png"
    add_image(s, fig, Inches(0.4), Inches(1.1), Inches(8.2))
    add_rect(s, Inches(8.9), Inches(1.1), Inches(4.0), Inches(2.1), RGBColor(0x1E, 0x3A, 0x5F))
    add_textbox(s, "← Poca solar (hoy)",
                Inches(9.1), Inches(1.2), Inches(3.6), Inches(0.45),
                font_size=13, color=TEXT_MUTE)
    add_textbox(s, "C1 lidera\n(96% cobertura)",
                Inches(9.1), Inches(1.65), Inches(3.6), Inches(0.55),
                font_size=14, color=BLUE, bold=True)
    add_rect(s, Inches(8.9), Inches(3.4), Inches(4.0), Inches(2.1), GREEN_BG)
    add_textbox(s, "Mucha solar (2030) →",
                Inches(9.1), Inches(3.5), Inches(3.6), Inches(0.45),
                font_size=13, color=TEXT_MUTE)
    add_textbox(s, "P2P domina\n(desde 110%)",
                Inches(9.1), Inches(3.95), Inches(3.6), Inches(0.55),
                font_size=14, color=GREEN, bold=True)
    add_rect(s, Inches(8.9), Inches(5.7), Inches(4.0), Inches(0.9), DARK_BOX)
    add_textbox(s, "Punto de cruce: φ ≈ 1.03 (103% cobertura)",
                Inches(9.1), Inches(5.75), Inches(3.7), Inches(0.8),
                font_size=13, color=YELLOW, bold=True)
    add_rect(s, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.75), DARK_BOX)
    add_textbox(
        s, "El Plan UPME 2025–2039 proyecta esa cobertura en 2030 — en 4 años.",
        Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.55),
        font_size=16, color=YELLOW, bold=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "Este es el hallazgo más importante de todo el trabajo. Hoy, con la "
        "cobertura solar que tienen estas cinco instituciones, la regulación actual "
        "C1 gana — por poquito. Pero hay un punto de quiebre: en cuanto la cobertura "
        "solar llega al 103%, P2P toma la delantera. Y desde ahí nunca vuelve a "
        "perder. ¿Por qué importa esto? Porque el Plan de Expansión dice que en 2030 "
        "vamos a tener esa cobertura y más. Esto no es un resultado académico lejano "
        "— es lo que va a pasar en los próximos cuatro años si Colombia no ajusta "
        "su regulación."
    )
    return s


def slide_13_robustez(prs):
    """Robustez del resultado."""
    s = new_blank_slide(prs)
    add_title(s, "El resultado es robusto — se mantiene en todos los escenarios probados", color=TEXT_PRI)
    fig = FIGURES / "fig_audit_calibration_robustness.png"
    add_image(s, fig, Inches(0.4), Inches(1.1), Inches(7.2))
    items = [
        (GREEN, "✓ Dos métodos de distribución",
         "Resultado estable con los dos métodos que permite la CREG 101 072 para repartir excedentes."),
        (GREEN, "✓ Todos los factores φ probados",
         "P2P domina en φ ∈ {1.1, 1.2, ..., 3.0} — de 110% a 288% de cobertura solar."),
        (GREEN, "✓ Ventaja monótona y estable",
         "La ventaja de P2P crece o se mantiene — no es un pico aislado ni un caso especial."),
        (GREEN, "✓ Mayoría individual",
         "4 de 5 instituciones prefieren P2P individualmente — no solo en el agregado."),
    ]
    for i, (col, titulo, desc) in enumerate(items):
        y = Inches(1.2) + i * Inches(1.3)
        add_rect(s, Inches(7.9), y, Inches(5.0), Inches(1.1), DARK_BOX)
        add_textbox(s, titulo, Inches(8.1), y + Inches(0.05), Inches(4.6), Inches(0.5),
                    font_size=14, color=col, bold=True)
        add_textbox(s, desc, Inches(8.1), y + Inches(0.55), Inches(4.6), Inches(0.5),
                    font_size=11, color=TEXT_SEC)
    add_speaker_note(s,
        "Una pregunta válida es: ¿esto solo aplica para agosto 2025 con estos parámetros "
        "exactos? La respuesta es no. El resultado se mantiene con los dos métodos "
        "permitidos por la regulación para distribuir excedentes. P2P domina en todos "
        "los factores solares probados, desde 110% hasta 288% de cobertura. La ventaja "
        "es monótona — no es un accidente estadístico. Y cuatro de las cinco "
        "instituciones lo prefieren de manera individual."
    )
    return s


# ──────────────────────────────────────────────
# BLOQUE 4 — PREGUNTA 3 (slides 14–16)
# ──────────────────────────────────────────────

def slide_14_implicacion_regulatoria(prs):
    """Implicación regulatoria."""
    s = new_blank_slide(prs)
    add_title(s, "¿Qué significa esto para Colombia?", color=TEXT_PRI)
    bloques = [
        (YELLOW, "Lo que muestra el estudio:",
         "La liquidación mensual de C1 y C2 no captura el valor hora a hora. "
         "Al promediar el mes, se pierde la diferencia entre mediodía con sol y la madrugada sin demanda."),
        (BLUE, "Lo que sugiere el resultado:",
         "A medida que Colombia escala la generación solar, los marcos regulatorios "
         "necesitarán mecanismos de liquidación horaria para maximizar el valor de esa energía."),
        (GREEN, "Lo que NO dice el estudio:",
         "No propone derogar la CREG 174. Muestra dónde está el techo y qué se gana "
         "al superarlo. La regulación actual funciona bien para coberturas bajas."),
    ]
    for i, (col, titulo, desc) in enumerate(bloques):
        y = Inches(1.4) + i * Inches(1.7)
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(1.5), RGBColor(0x1E, 0x29, 0x3B))
        add_rect(s, Inches(0.4), y, Inches(0.12), Inches(1.5), col)
        add_textbox(s, titulo, Inches(0.7), y + Inches(0.1),
                    Inches(3.5), Inches(0.5), font_size=15, color=col, bold=True)
        add_textbox(s, desc, Inches(0.7), y + Inches(0.6),
                    Inches(12.0), Inches(0.8), font_size=13, color=TEXT_SEC)
    add_speaker_note(s,
        "¿Qué significa todo esto para Colombia? No significa que la regulación actual "
        "esté mal — significa que tiene un techo. Cuando la generación solar era poca, "
        "esa regulación era suficiente. Pero a medida que instalamos más paneles, el "
        "mecanismo de liquidar todo al final del mes deja de capturar el valor real de "
        "esa energía. Lo que este trabajo dice es: si vamos a cumplir el Plan de "
        "Expansión, vale la pena pensar en cómo liquidar la energía hora por hora."
    )
    return s


def slide_15_limites(prs):
    """Límites honestos del estudio."""
    s = new_blank_slide(prs)
    add_title(s, "Límites del estudio — qué no cubre este trabajo todavía", color=ORANGE)
    limites = [
        "1 mes de datos (agosto 2025) en 5 instituciones de Pasto. "
        "Pueden existir variaciones estacionales o en otras ciudades.",
        "Sin respuesta de demanda (α = 0): no se modeló ajuste de consumo según precios P2P. "
        "El beneficio real podría ser mayor con DR activada.",
        "Tarifas homogeneizadas al perfil comercial. "
        "El análisis con tarifas mixtas queda para la tesis completa.",
        "Costos de transacción de implementar P2P no modelados: "
        "plataforma, medición horaria, infraestructura de comunicación.",
    ]
    for i, lim in enumerate(limites):
        y = Inches(1.4) + i * Inches(1.3)
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(1.1), DARK_BOX)
        add_textbox(s, "⚠", Inches(0.6), y + Inches(0.2), Inches(0.5), Inches(0.6),
                    font_size=22, color=ORANGE)
        add_textbox(s, lim, Inches(1.3), y + Inches(0.2), Inches(11.4), Inches(0.8),
                    font_size=13, color=TEXT_SEC)
    add_rect(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.7), DARK_BOX)
    add_textbox(
        s, "Los límites acotan la conclusión — no la anulan. "
           "Ser honestos con ellos fortalece, no debilita, los resultados que sí aplican.",
        Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6),
        font_size=13, color=TEXT_MUTE, italic=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "Quiero ser honesto con lo que este estudio no cubre. Un mes de datos de agosto. "
        "Cinco instituciones. Sin modelar que los usuarios ajusten su consumo según los "
        "precios P2P. Esos son los límites reales. Y los digo porque si los ocultara, "
        "alguien los encontraría y pondría en duda todo lo demás. Los límites no anulan "
        "la conclusión — la acotan, que es lo correcto en ciencia."
    )
    return s


def slide_16_proximos_pasos(prs):
    """Próximos pasos MTE y tesis."""
    s = new_blank_slide(prs)
    add_title(s, "¿Qué sigue? Próximos pasos de la tesis y del MTE", color=TEXT_PRI)
    add_rect(s, Inches(0.4), Inches(1.3), Inches(6.0), Inches(4.6), DARK_BOX)
    add_textbox(s, "Tesis (próximos meses)",
                Inches(0.6), Inches(1.45), Inches(5.6), Inches(0.55),
                font_size=18, color=PURPLE, bold=True)
    tesis_items = [
        "Horizonte completo 5160 h (Jul 2025–Ene 2026)",
        "Análisis de sensibilidad global Sobol",
        "Incorporar respuesta de demanda (α > 0)",
        "Tarifas heterogéneas por institución",
    ]
    for i, item in enumerate(tesis_items):
        add_textbox(s, f"→  {item}",
                    Inches(0.6), Inches(2.1) + i * Inches(0.85), Inches(5.6), Inches(0.7),
                    font_size=14, color=TEXT_SEC)
    add_rect(s, Inches(6.9), Inches(1.3), Inches(6.0), Inches(4.6), DARK_BOX)
    add_textbox(s, "MTE (infraestructura)",
                Inches(7.1), Inches(1.45), Inches(5.6), Inches(0.55),
                font_size=18, color=BLUE, bold=True)
    mte_items = [
        "Datos Jul 2025–Ene 2026 disponibles",
        "Validación estacional en curso",
        "Base empírica para estudios futuros",
        "Paper WEEF 2026 enviado — sept. 22–24",
    ]
    for i, item in enumerate(mte_items):
        add_textbox(s, f"→  {item}",
                    Inches(7.1), Inches(2.1) + i * Inches(0.85), Inches(5.6), Inches(0.7),
                    font_size=14, color=TEXT_SEC)
    add_rect(s, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.65), DARK_BOX)
    add_textbox(
        s, "El MTE no termina aquí — es la base empírica de toda la tesis.",
        Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.55),
        font_size=15, color=YELLOW, bold=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s,
        "Lo que sigue: el MTE seguirá midiendo hasta enero de 2026. Eso nos dará un "
        "horizonte completo de 5160 horas para confirmar si lo que encontramos en "
        "agosto se mantiene en todas las estaciones. Y el paper fue enviado a "
        "WEEF 2026 — lo presentamos en septiembre."
    )
    return s


# ──────────────────────────────────────────────
# BLOQUE 5 — CIERRE (slides 17–18)
# ──────────────────────────────────────────────

def slide_17_conclusion(prs):
    """Conclusión — 3 bullets."""
    s = new_blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), W, Inches(0.08), YELLOW)
    add_title(s, "Lo que nos dicen los datos de Pasto", color=YELLOW, size=32)
    conclusiones = [
        (GREEN,  "①",
         "El MTE produce evidencia real.",
         "744 horas en 5 instituciones de Pasto — la base que Colombia necesita "
         "para tomar decisiones sobre energía comunitaria."),
        (BLUE,   "②",
         "P2P supera a la regulación actual desde el 103% de cobertura.",
         "4 de 5 instituciones lo prefieren individualmente. "
         "La ventaja es robusta y monótona."),
        (ORANGE, "③",
         "Con el plan solar 2030, la regulación necesitará liquidación horaria.",
         "Colombia tiene tiempo de adaptarse — el punto de cruce llega en 4 años."),
    ]
    for i, (col, num, titulo, desc) in enumerate(conclusiones):
        y = Inches(1.4) + i * Inches(1.7)
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(1.5), DARK_BOX)
        add_textbox(s, num, Inches(0.55), y + Inches(0.3), Inches(0.6), Inches(0.8),
                    font_size=28, color=col, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, titulo, Inches(1.3), y + Inches(0.1), Inches(11.3), Inches(0.55),
                    font_size=17, color=col, bold=True)
        add_textbox(s, desc, Inches(1.3), y + Inches(0.7), Inches(11.3), Inches(0.7),
                    font_size=13, color=TEXT_SEC)
    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), YELLOW)
    add_speaker_note(s,
        "Tres ideas para llevarse. El MTE produce datos reales que Colombia necesita. "
        "P2P gana desde el 103% de cobertura solar, con 4 de 5 instituciones a favor. "
        "Y con el plan solar del Gobierno, esa cobertura llega en 2030 — la regulación "
        "tiene tiempo de adaptarse."
    )
    return s


def slide_18_preguntas(prs):
    """Slide de preguntas."""
    s = new_blank_slide(prs)
    add_textbox(s, "¿Preguntas?",
                Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
                font_size=52, color=TEXT_PRI, bold=True, align=PP_ALIGN.CENTER)
    contactos = [
        (BLUE,   "📧 Datos MTE",          "bralopez@udenar.edu.co"),
        (PURPLE, "💻 Código y resultados", "github.com / SistemaBL"),
        (GREEN,  "📄 Paper completo",      "WEEF 2026 · Sept 22–24"),
    ]
    for i, (col, label, val) in enumerate(contactos):
        x = Inches(0.6) + i * Inches(4.2)
        add_rect(s, x, Inches(3.3), Inches(3.9), Inches(1.4), DARK_BOX)
        add_textbox(s, label, x + Inches(0.15), Inches(3.45), Inches(3.6), Inches(0.5),
                    font_size=14, color=col, bold=True)
        add_textbox(s, val, x + Inches(0.15), Inches(3.95), Inches(3.6), Inches(0.6),
                    font_size=13, color=TEXT_SEC)
    add_textbox(
        s, "Slides de respaldo disponibles: metodología detallada · tablas completas · ecuaciones del modelo",
        Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.5),
        font_size=12, color=TEXT_MUTE, italic=True, align=PP_ALIGN.CENTER
    )
    add_speaker_note(s, "Slide de preguntas. Referir a slides de respaldo según el tipo de pregunta.")
    return s


# ──────────────────────────────────────────────
# BACKUP SLIDES
# ──────────────────────────────────────────────

def slide_b1_ecuaciones(prs):
    """Backup: ecuaciones del modelo."""
    s = new_blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), W, Inches(0.35), RGBColor(0x33, 0x41, 0x55))
    add_textbox(s, "BACKUP — Metodología: Ecuaciones del modelo",
                Inches(0.5), Inches(0.05), Inches(12.0), Inches(0.25),
                font_size=11, color=TEXT_MUTE, italic=True)
    add_title(s, "Dinámica de Replicador + Stackelberg: formulación completa", color=PURPLE)
    add_textbox(
        s, "Problema del vendedor j (líder Stackelberg):",
        Inches(0.5), Inches(1.4), Inches(12.0), Inches(0.45),
        font_size=14, color=TEXT_SEC, bold=True
    )
    add_textbox(
        s, "max  Σᵢ πⱼ Pⱼᵢ − Hⱼ(Σᵢ Pⱼᵢ)     s.t.  0 ≤ Pⱼᵢ ≤ Gⱼ − Dⱼ (∀i)",
        Inches(0.7), Inches(1.9), Inches(11.6), Inches(0.55),
        font_size=16, color=YELLOW
    )
    add_textbox(
        s, "Hⱼ(P) = aⱼP² + bⱼP + cⱼ   con  aⱼ=0,  bⱼ≈241 COP/kWh (PV puro),  cⱼ=0",
        Inches(0.7), Inches(2.5), Inches(11.6), Inches(0.45),
        font_size=13, color=TEXT_SEC
    )
    add_textbox(
        s, "Dinámica del comprador i (seguidor):",
        Inches(0.5), Inches(3.2), Inches(12.0), Inches(0.45),
        font_size=14, color=TEXT_SEC, bold=True
    )
    add_textbox(
        s, "ẋᵢⱼ = xᵢⱼ (fᵢⱼ − f̄ᵢ)     fᵢⱼ = (πgs,i − πⱼ)(πⱼ − πgb)",
        Inches(0.7), Inches(3.7), Inches(11.6), Inches(0.55),
        font_size=16, color=YELLOW
    )
    add_textbox(
        s, "xᵢⱼ: fracción de demanda del comprador i asignada al vendedor j\n"
           "πgs,i: precio de venta a red (techo)   πgb: precio de bolsa (piso)\n"
           "Convergencia: tolerancia 10⁻³, EWMA τ=10⁻³, 2 iteraciones Stackelberg",
        Inches(0.7), Inches(4.35), Inches(11.6), Inches(1.1),
        font_size=12, color=TEXT_MUTE
    )
    add_speaker_note(s, "Backup para preguntas técnicas sobre la formulación matemática.")
    return s


def slide_b2_calibracion(prs):
    """Backup: tabla de calibración de parámetros."""
    s = new_blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), W, Inches(0.35), RGBColor(0x33, 0x41, 0x55))
    add_textbox(s, "BACKUP — Calibración de parámetros",
                Inches(0.5), Inches(0.05), Inches(12.0), Inches(0.25),
                font_size=11, color=TEXT_MUTE, italic=True)
    add_title(s, "Parámetros del modelo: fuente, valor y invariancia", color=BLUE)
    params = [
        ("aⱼ", "0",              "PV puro — sin coste cuadrático (IRENA 2024, Chacón Tabla I)"),
        ("bⱼ", "225–241 COP/kWh","LCOE solar UPME Colombia; Fronius=241, Cesmag=225 (−6.7%)"),
        ("cⱼ", "0",              "Coste fijo cero para renovables (invariante analítico bajo α=0)"),
        ("λⱼ", "100",            "Invariante analítico bajo α=0 — no afecta el equilibrio (P*, π*)"),
        ("θⱼ", "0.5",            "Invariante analítico bajo α=0 — no afecta el equilibrio (P*, π*)"),
        ("ηᵢ", "0.1",            "Invariante analítico bajo α=0 — no afecta el equilibrio (P*, π*)"),
        ("α",  "0",              "Sin respuesta de demanda — aísla el mecanismo P2P"),
    ]
    for i, (param, val, fuente) in enumerate(params):
        y = Inches(1.4) + i * Inches(0.72)
        bg = DARK_BOX if i % 2 == 0 else RGBColor(0x0F, 0x17, 0x2A)
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.62), bg)
        add_textbox(s, param, Inches(0.55), y + Inches(0.1), Inches(0.7), Inches(0.45),
                    font_size=14, color=YELLOW, bold=True)
        add_textbox(s, val, Inches(1.4), y + Inches(0.1), Inches(2.5), Inches(0.45),
                    font_size=13, color=GREEN, bold=True)
        add_textbox(s, fuente, Inches(4.1), y + Inches(0.1), Inches(8.7), Inches(0.45),
                    font_size=11, color=TEXT_SEC)
    add_speaker_note(s, "Backup para preguntas sobre calibración de parámetros del modelo.")
    return s


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # Bloque 1 — Gancho
    slide_01_portada(prs)
    slide_02_pregunta_provocadora(prs)
    slide_03_mte_intro(prs)
    # Bloque 2 — Pregunta 1
    slide_04_que_es_mte(prs)
    slide_05_perfiles(prs)
    slide_06_problema_fondo(prs)
    # Bloque 3 — Pregunta 2
    slide_07_que_es_p2p(prs)
    slide_08_algoritmo(prs)
    slide_09_esquemas_colombianos(prs)
    slide_10_resultado_principal(prs)
    slide_11_por_agente(prs)
    slide_12_transicion_fase(prs)
    slide_13_robustez(prs)
    # Bloque 4 — Pregunta 3
    slide_14_implicacion_regulatoria(prs)
    slide_15_limites(prs)
    slide_16_proximos_pasos(prs)
    # Bloque 5 — Cierre
    slide_17_conclusion(prs)
    slide_18_preguntas(prs)
    # Backup slides
    slide_b1_ecuaciones(prs)
    slide_b2_calibracion(prs)

    out = Path(__file__).parent / "MTE_Avances_2026.pptx"
    prs.save(out)
    print(f"Guardado: {out}  ({len(prs.slides)} slides)")
